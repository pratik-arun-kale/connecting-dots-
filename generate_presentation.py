"""
Generates GenAI_Project_Presentation.pptx — 5-slide deck for Context Workspace + Advanced RAG.
Run: python generate_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x0F, 0x1A)   # near-black navy
ACCENT    = RGBColor(0x6E, 0x56, 0xCF)   # indigo-violet
ACCENT2   = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
GREEN     = RGBColor(0x34, 0xD3, 0x99)   # emerald
AMBER     = RGBColor(0xFB, 0xBF, 0x24)   # amber
RED_SOFT  = RGBColor(0xF8, 0x71, 0x71)   # soft red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xA0, 0xAE, 0xC0)
DARKGRAY  = RGBColor(0x1E, 0x24, 0x3B)   # card background
LIGHTGRAY = RGBColor(0xC8, 0xD3, 0xE8)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helper utilities ──────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if False else 1,  # 1 = MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox


def add_label_value(slide, label, value, left, top, width,
                    lsize=10, vsize=13, lcolor=GRAY, vcolor=WHITE):
    add_text_box(slide, label, left, top, width, Inches(0.25),
                 font_size=lsize, color=lcolor, bold=True)
    add_text_box(slide, value, left, top + Inches(0.22), width, Inches(0.35),
                 font_size=vsize, color=vcolor)


def add_pill(slide, text, left, top, width, height, bg, txt_color=WHITE, size=9):
    r = add_rect(slide, left, top, width, height, fill_color=bg)
    # rounded corners via XML
    sp = r._element
    sp_pr = sp.find('{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}spPr') or \
            sp.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spPr') or \
            sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}spPr')
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = txt_color
    run.font.name = "Segoe UI"
    return r


def add_divider(slide, top, color=ACCENT, width_pct=0.9):
    left = SLIDE_W * (1 - width_pct) / 2
    w    = SLIDE_W * width_pct
    line = slide.shapes.add_shape(1, left, top, w, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def section_header(slide, label, top):
    add_text_box(slide, label.upper(), Inches(0.5), top, Inches(12), Inches(0.3),
                 font_size=8, bold=True, color=ACCENT, italic=False)
    add_divider(slide, top + Inches(0.28), color=ACCENT, width_pct=0.94)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — GenAI / Agentic AI Project Experience
# ══════════════════════════════════════════════════════════════════════════════

s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, BG)

# Left accent bar
add_rect(s1, 0, 0, Inches(0.06), SLIDE_H, fill_color=ACCENT)

# Slide number badge
add_rect(s1, Inches(12.6), Inches(0.12), Inches(0.55), Inches(0.3), fill_color=ACCENT)
add_text_box(s1, "01 / 05", Inches(12.58), Inches(0.1), Inches(0.6), Inches(0.35),
             font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title block
add_text_box(s1, "GenAI / Agentic AI", Inches(0.5), Inches(0.2), Inches(10), Inches(0.65),
             font_size=36, bold=True, color=WHITE)
add_text_box(s1, "Project Experience", Inches(0.5), Inches(0.78), Inches(10), Inches(0.55),
             font_size=36, bold=True, color=ACCENT2)
add_text_box(s1, "Context Workspace + Advanced RAG Pipeline", Inches(0.5), Inches(1.28), Inches(10), Inches(0.3),
             font_size=13, bold=False, color=GRAY)

add_divider(s1, Inches(1.62), color=ACCENT)

# ── Project summary cards (row) ───────────────────────────────────────────────
cards = [
    ("Chrome Extension",  "Manifest V3 · React 18\nCaptures AI conversations\nfrom ChatGPT, Claude,\nGemini, Perplexity",  ACCENT),
    ("FastAPI Backend",   "Async PostgreSQL · SQLAlchemy\nContext ingestion API\nIdempotent capture\nStructured logging",    ACCENT2),
    ("Next.js Dashboard", "TanStack Query · TypeScript\nProject management\nContext viewer\nReal-time sync",               GREEN),
    ("Advanced RAG",      "Hybrid retrieval · BM25+Vector\nCross-encoder reranking\nExtractiveQA · ChromaDB\nStreamlit UI", AMBER),
]

card_w = Inches(2.95)
card_h = Inches(2.2)
gap    = Inches(0.18)
start_x = Inches(0.35)
start_y = Inches(1.78)

for i, (title, body, color) in enumerate(cards):
    cx = start_x + i * (card_w + gap)
    add_rect(s1, cx, start_y, card_w, card_h, fill_color=DARKGRAY, line_color=color, line_width=Pt(1.5))
    # Color top stripe
    add_rect(s1, cx, start_y, card_w, Inches(0.06), fill_color=color)
    add_text_box(s1, title, cx + Inches(0.15), start_y + Inches(0.1), card_w - Inches(0.2), Inches(0.35),
                 font_size=13, bold=True, color=color)
    add_text_box(s1, body, cx + Inches(0.15), start_y + Inches(0.45), card_w - Inches(0.2), Inches(1.7),
                 font_size=10, color=LIGHTGRAY)

# ── RAG Pipeline flow ─────────────────────────────────────────────────────────
section_header(s1, "Advanced RAG Pipeline (detailed)", Inches(4.13))

stages = [
    ("TXT Docs\n9 HR policies", BG),
    ("Chunker\n250 tok / 50 overlap", DARKGRAY),
    ("Embeddings\nall-MiniLM-L6-v2", DARKGRAY),
    ("ChromaDB\nVector Store", DARKGRAY),
    ("BM25 +\nVector Retrieval", DARKGRAY),
    ("RRF Fusion\nk = 60", DARKGRAY),
    ("Cross-Encoder\nReranker", DARKGRAY),
    ("Confidence\nGrader HIGH/MED/LOW", DARKGRAY),
    ("RoBERTa\nExtractiveQA", DARKGRAY),
    ("Answer +\nCitation", GREEN),
]

box_w = Inches(1.15)
box_h = Inches(0.72)
y_flow = Inches(4.48)
x_flow = Inches(0.25)
x_step = (SLIDE_W - Inches(0.5)) / len(stages)

for i, (label, fill) in enumerate(stages):
    cx = x_flow + i * x_step
    fill_c = fill if fill != BG else ACCENT
    border_c = ACCENT if fill == DARKGRAY else fill_c
    add_rect(s1, cx, y_flow, box_w, box_h, fill_color=fill_c, line_color=border_c, line_width=Pt(1))
    add_text_box(s1, label, cx + Inches(0.04), y_flow + Inches(0.05),
                 box_w - Inches(0.05), box_h, font_size=7.5, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow
    if i < len(stages) - 1:
        arr_x = cx + box_w + Inches(0.01)
        add_text_box(s1, "→", arr_x, y_flow + Inches(0.22), Inches(0.22), Inches(0.3),
                     font_size=12, color=ACCENT2, align=PP_ALIGN.CENTER)

# Corrective RAG note
add_rect(s1, Inches(0.25), Inches(5.32), Inches(8.5), Inches(0.45),
         fill_color=RGBColor(0x1a, 0x1f, 0x35), line_color=AMBER, line_width=Pt(1))
add_text_box(s1, "⚡ Corrective RAG: If confidence is LOW/MEDIUM → synonym expansion or query simplification → re-retrieve → re-grade (max 2 attempts)",
             Inches(0.35), Inches(5.35), Inches(8.3), Inches(0.4), font_size=8.5, color=AMBER)

# Context Capture pipeline
add_rect(s1, Inches(9.0), Inches(5.25), Inches(4.1), Inches(1.55),
         fill_color=RGBColor(0x0e, 0x18, 0x30), line_color=ACCENT, line_width=Pt(1))
add_text_box(s1, "Context Capture Pipeline", Inches(9.1), Inches(5.28), Inches(3.9), Inches(0.3),
             font_size=10, bold=True, color=ACCENT2)

cap_steps = ["AI Tab (ChatGPT/Claude/Gemini/Perplexity)",
             "executeScript → DOM extraction → JSON",
             "SHA-256 idempotency key",
             "POST /api/v1/projects/{id}/capture",
             "PostgreSQL → Next.js dashboard"]

for j, step in enumerate(cap_steps):
    add_text_box(s1, f"  {j+1}. {step}", Inches(9.1), Inches(5.58) + j * Inches(0.2),
                 Inches(3.9), Inches(0.22), font_size=7.5, color=LIGHTGRAY)

# Speaker notes
notes_1 = s1.notes_slide.notes_text_frame
notes_1.text = (
    "SLIDE 1 — GenAI Project Experience\n\n"
    "This project is a full-stack AI productivity tool called 'Context Workspace.' "
    "It has four components: a Chrome Extension for capturing AI conversations, "
    "a FastAPI async backend persisting data to PostgreSQL, a Next.js dashboard for reviewing captured context, "
    "and a standalone Advanced RAG pipeline for HR policy Q&A.\n\n"
    "The RAG pipeline uses a full advanced retrieval chain: "
    "token-based chunking → sentence-transformer embeddings → ChromaDB vector store → "
    "hybrid BM25 + vector retrieval → RRF fusion → cross-encoder reranking → confidence grading → "
    "corrective retrieval fallback → RoBERTa extractive QA → answer with citation.\n\n"
    "Key differentiator: confidence grading automatically classifies retrieval quality as HIGH/MEDIUM/LOW "
    "based on cross-encoder score and retrieval agreement between BM25 and vector search. "
    "LOW confidence triggers corrective retrieval (synonym expansion or query simplification)."
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Solution Architecture
# ══════════════════════════════════════════════════════════════════════════════

s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, BG)
add_rect(s2, 0, 0, Inches(0.06), SLIDE_H, fill_color=ACCENT2)

add_rect(s2, Inches(12.6), Inches(0.12), Inches(0.55), Inches(0.3), fill_color=ACCENT2)
add_text_box(s2, "02 / 05", Inches(12.58), Inches(0.1), Inches(0.6), Inches(0.35),
             font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(s2, "Solution Architecture", Inches(0.5), Inches(0.12), Inches(10), Inches(0.55),
             font_size=32, bold=True, color=WHITE)
add_text_box(s2, "Context Workspace — full system diagram", Inches(0.5), Inches(0.62), Inches(10), Inches(0.28),
             font_size=12, color=GRAY)

add_divider(s2, Inches(0.95), color=ACCENT2)

# ── LAYER LABELS (left column) ────────────────────────────────────────────────
layers = [
    ("USER LAYER",       Inches(1.12), GRAY),
    ("CAPTURE LAYER",    Inches(2.05), GRAY),
    ("PROCESSING LAYER", Inches(3.0),  GRAY),
    ("STORAGE LAYER",    Inches(4.0),  GRAY),
    ("PRESENTATION",     Inches(5.0),  GRAY),
    ("RAG ENGINE",       Inches(5.95), GRAY),
]

for label, y, color in layers:
    add_text_box(s2, label, Inches(0.08), y, Inches(1.15), Inches(0.3),
                 font_size=7, bold=True, color=color, align=PP_ALIGN.RIGHT)

# ── ARCHITECTURE BOXES ────────────────────────────────────────────────────────

def arch_box(slide, label, sublabel, left, top, width, height, fill, border, tcolor=WHITE, subcolor=GRAY):
    add_rect(slide, left, top, width, height, fill_color=fill, line_color=border, line_width=Pt(1.5))
    add_text_box(slide, label, left + Inches(0.08), top + Inches(0.06),
                 width - Inches(0.12), Inches(0.28), font_size=10, bold=True, color=tcolor)
    if sublabel:
        add_text_box(slide, sublabel, left + Inches(0.08), top + Inches(0.3),
                     width - Inches(0.12), height - Inches(0.3), font_size=8, color=subcolor)


# User Layer
arch_box(s2, "User / Browser", "ChatGPT · Claude · Gemini · Perplexity",
         Inches(1.35), Inches(1.05), Inches(2.8), Inches(0.7), DARKGRAY, GRAY)

# Chrome Extension
arch_box(s2, "Chrome Extension (MV3)", "React 18 · Tailwind · Framer Motion\nZustand · Popup + Side Panel",
         Inches(4.4), Inches(1.05), Inches(3.0), Inches(0.7), DARKGRAY, ACCENT)

# Background Service Worker
arch_box(s2, "Background Service Worker", "executeScript DOM injection\nCaptureQueue (durable · retry)\nSessionOrchestrator · TabManager",
         Inches(4.4), Inches(1.95), Inches(3.0), Inches(0.82), DARKGRAY, ACCENT)

# FastAPI Backend
arch_box(s2, "FastAPI Backend (Python 3.12)", "ExtensionCORSMiddleware (pure-ASGI)\nRoutes: /projects /sessions /contexts /capture\nService → Repository pattern",
         Inches(7.75), Inches(1.95), Inches(3.15), Inches(0.82), DARKGRAY, ACCENT2)

# PostgreSQL
arch_box(s2, "PostgreSQL", "Projects · Sessions · Contexts\nAlembic migrations (5 versions)\nJSONB for raw_content",
         Inches(7.75), Inches(3.0), Inches(1.5), Inches(0.7), DARKGRAY, GREEN)

# Redis
arch_box(s2, "Redis", "Future: queue\n& caching layer",
         Inches(9.45), Inches(3.0), Inches(1.45), Inches(0.7), DARKGRAY, AMBER)

# Next.js Dashboard
arch_box(s2, "Next.js Dashboard", "TanStack Query · TypeScript\nProject view · Context list\nPlatform badge · Message viewer",
         Inches(4.4), Inches(3.0), Inches(3.0), Inches(0.7), DARKGRAY, ACCENT2)

# Content Scripts
arch_box(s2, "Content Scripts", "chatgpt-cs · claude-cs\ngemini-cs · perplexity-cs\nAuto-scroll · DOM extraction",
         Inches(1.35), Inches(2.87), Inches(2.8), Inches(0.7), DARKGRAY, ACCENT)

# Advanced RAG (bottom)
add_rect(s2, Inches(1.35), Inches(5.82), Inches(9.55), Inches(0.06), fill_color=AMBER)
arch_box(s2, "Advanced RAG Pipeline (Streamlit UI)",
         "BM25Retriever (rank-bm25) + VectorRetriever (ChromaDB + all-MiniLM-L6-v2) → RRF Fusion → CrossEncoder Reranker → Confidence Grader → CorrectiveRAG → RoBERTa ExtractiveQA",
         Inches(1.35), Inches(5.9), Inches(9.55), Inches(0.82), RGBColor(0x18, 0x14, 0x0a), AMBER, tcolor=AMBER, subcolor=LIGHTGRAY)

# ── CONNECTORS (text arrows) ──────────────────────────────────────────────────
def arrow(slide, x, y, w=Inches(0.4), h=Inches(0.28), label="→", color=ACCENT2):
    add_text_box(slide, label, x, y, w, h, font_size=14, color=color, align=PP_ALIGN.CENTER, bold=True)

# User → Extension
arrow(s2, Inches(4.12), Inches(1.28), label="→")
# Extension → SW
add_text_box(s2, "↓", Inches(5.8), Inches(1.73), Inches(0.4), Inches(0.26), font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
# CS → Extension
arrow(s2, Inches(4.12), Inches(3.1), label="←")
# SW → Backend
arrow(s2, Inches(7.37), Inches(2.25), label="→")
# Backend → PG
add_text_box(s2, "↓", Inches(8.28), Inches(2.77), Inches(0.4), Inches(0.26), font_size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
# Backend → Redis
add_text_box(s2, "↓", Inches(10.1), Inches(2.77), Inches(0.4), Inches(0.26), font_size=14, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
# PG → Dashboard
arrow(s2, Inches(7.37), Inches(3.28), label="←")

# Tech-stack pills row
pills = [
    ("Python 3.12", ACCENT), ("FastAPI", ACCENT2), ("SQLAlchemy async", GREEN), ("PostgreSQL", RGBColor(0x33,0x6B,0x9E)),
    ("ChromaDB", AMBER), ("sentence-transformers", RGBColor(0xEC,0x48,0x99)), ("React 18", ACCENT2), ("Next.js", WHITE),
]
pill_y = Inches(4.88)
px = Inches(1.35)
for text, color in pills:
    pw = Inches(len(text) * 0.085 + 0.25)
    add_rect(s2, px, pill_y, pw, Inches(0.28), fill_color=RGBColor(0x1e,0x24,0x3b), line_color=color, line_width=Pt(1))
    add_text_box(s2, text, px + Inches(0.05), pill_y, pw, Inches(0.28), font_size=8.5, color=color, bold=True)
    px += pw + Inches(0.12)

notes_2 = s2.notes_slide.notes_text_frame
notes_2.text = (
    "SLIDE 2 — Solution Architecture\n\n"
    "The system has five layers:\n"
    "1. User Layer: The user browses AI platforms (ChatGPT, Claude, Gemini, Perplexity)\n"
    "2. Capture Layer: Chrome Extension MV3 with React popup; Content Scripts auto-scroll and extract DOM\n"
    "3. Processing Layer: Background Service Worker uses chrome.scripting.executeScript to inject extraction logic, "
    "builds idempotency keys via SHA-256, enqueues captures with exponential backoff retry\n"
    "4. Storage Layer: FastAPI backend with pure-ASGI ExtensionCORSMiddleware shim, "
    "PostgreSQL via SQLAlchemy async, Redis (future)\n"
    "5. Presentation: Next.js dashboard with TanStack Query\n\n"
    "The Advanced RAG Pipeline is a separate module that will integrate with the backend "
    "to provide intelligent Q&A over captured context. It uses hybrid BM25+vector retrieval, "
    "RRF fusion, cross-encoder reranking, and RoBERTa extractive QA."
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Production Issue & Resolution
# ══════════════════════════════════════════════════════════════════════════════

s3 = prs.slides.add_slide(blank_layout)
set_bg(s3, BG)
add_rect(s3, 0, 0, Inches(0.06), SLIDE_H, fill_color=RED_SOFT)

add_rect(s3, Inches(12.6), Inches(0.12), Inches(0.55), Inches(0.3), fill_color=RED_SOFT)
add_text_box(s3, "03 / 05", Inches(12.58), Inches(0.1), Inches(0.6), Inches(0.35),
             font_size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)

add_text_box(s3, "Production Issue & Resolution", Inches(0.5), Inches(0.12), Inches(11), Inches(0.55),
             font_size=32, bold=True, color=WHITE)
add_text_box(s3, "Engineering challenges diagnosed and resolved during development",
             Inches(0.5), Inches(0.62), Inches(11), Inches(0.28), font_size=12, color=GRAY)
add_divider(s3, Inches(0.95), color=RED_SOFT)

# ── Issue cards ───────────────────────────────────────────────────────────────
issues = [
    {
        "title": "Issue 1: CORS Headers Silently Swallowed",
        "severity": "CRITICAL",
        "sev_color": RED_SOFT,
        "symptom": "Chrome extension fetch to http://localhost:8000 returned no CORS headers → all API calls blocked with 'No Access-Control-Allow-Origin' error.",
        "root_cause": "Triple-layer problem: (1) Starlette BaseHTTPMiddleware (RequestLoggingMiddleware) wraps the inner ASGI send() before CORSMiddleware can write headers. (2) allow_origin_regex was correct but headers got consumed. (3) Ghost uvicorn process from a previous session still held port 8000 — new (fixed) server couldn't bind, so old (unfixed) server kept responding.",
        "diagnosis": "Verified regex in isolation via TestClient → passed. Confirmed middleware order. Discovered ghost process via netstat + WMI Get-WmiObject Win32_Process.",
        "resolution": "Created ExtensionCORSMiddleware pure-ASGI shim — runs at the outermost ASGI layer (above all Starlette middleware), intercepts http.response.start and appends CORS headers directly. Killed ghost processes via PowerShell WMI query.",
        "learning": "Never use BaseHTTPMiddleware for CORS — it wraps ASGI send() and eats inner middleware headers. Pure-ASGI middleware avoids this. Always check for ghost processes when a server change appears to have no effect.",
    },
    {
        "title": "Issue 2: 'Could Not Extract Conversation' Error",
        "severity": "HIGH",
        "sev_color": AMBER,
        "symptom": "Clicking 'Capture Context' in popup always returned: 'Could not extract conversation. Try scrolling the page first.'",
        "root_cause": "chrome.tabs.sendMessage requires the target content script to be already registered and listening. If the content script hadn't fully initialized, or the active tab was on any non-AI page, the message had no recipient and threw a 'Could not establish connection' error.",
        "diagnosis": "Traced the capture flow: popup → background → chrome.tabs.sendMessage → content script (unreachable). Confirmed content scripts only register on specific domains.",
        "resolution": "Replaced chrome.tabs.sendMessage with chrome.scripting.executeScript — injects the full extraction function directly into the active tab at capture time. No pre-registration required. Added auto-scroll, platform auto-detection from hostname, structured message extraction, and 50k char body fallback.",
        "learning": "For ad-hoc DOM extraction in Manifest V3, executeScript is more reliable than sendMessage because it doesn't depend on content script lifecycle. Use sendMessage only for persistent background communication.",
    },
]

col_w = Inches(6.1)
col_gap = Inches(0.2)
cols = [Inches(0.25), Inches(6.55)]
heights = [Inches(2.55), Inches(2.55)]
tops = [Inches(1.05), Inches(1.05)]

for idx, (issue, col_x, col_top) in enumerate(zip(issues, cols, tops)):
    ch = heights[idx]
    add_rect(s3, col_x, col_top, col_w, ch,
             fill_color=DARKGRAY, line_color=issue["sev_color"], line_width=Pt(2))
    # Top severity bar
    add_rect(s3, col_x, col_top, col_w, Inches(0.07), fill_color=issue["sev_color"])

    # Title + badge
    add_text_box(s3, issue["title"], col_x + Inches(0.12), col_top + Inches(0.1),
                 col_w - Inches(0.5), Inches(0.3), font_size=12, bold=True, color=issue["sev_color"])
    # Severity badge
    bw = Inches(0.75)
    add_rect(s3, col_x + col_w - bw - Inches(0.08), col_top + Inches(0.1),
             bw, Inches(0.25), fill_color=issue["sev_color"])
    add_text_box(s3, issue["severity"],
                 col_x + col_w - bw - Inches(0.08), col_top + Inches(0.1),
                 bw, Inches(0.25), font_size=7.5, bold=True, color=BG, align=PP_ALIGN.CENTER)

    rows = [
        ("SYMPTOM",     issue["symptom"],     RED_SOFT),
        ("ROOT CAUSE",  issue["root_cause"],  AMBER),
        ("RESOLUTION",  issue["resolution"],  GREEN),
    ]
    row_y = col_top + Inches(0.42)
    for label, text, lc in rows:
        add_text_box(s3, label, col_x + Inches(0.12), row_y,
                     col_w - Inches(0.2), Inches(0.18), font_size=7.5, bold=True, color=lc)
        row_y += Inches(0.18)
        add_text_box(s3, text, col_x + Inches(0.12), row_y,
                     col_w - Inches(0.2), Inches(0.4), font_size=8.5, color=LIGHTGRAY)
        row_y += Inches(0.4)

# Learning row
add_rect(s3, Inches(0.25), Inches(3.7), Inches(12.85), Inches(0.5),
         fill_color=RGBColor(0x0e, 0x18, 0x30), line_color=GREEN, line_width=Pt(1))
add_text_box(s3, "📚  KEY LEARNINGS", Inches(0.35), Inches(3.72), Inches(2.0), Inches(0.46),
             font_size=9, bold=True, color=GREEN)
add_text_box(s3,
             "Issue 1: Never use BaseHTTPMiddleware for CORS in FastAPI — it silently consumes headers. Always verify no orphan processes own the port.  "
             "|  Issue 2: Use chrome.scripting.executeScript for reliable on-demand DOM extraction in MV3 extensions — no content script pre-registration required.",
             Inches(2.3), Inches(3.72), Inches(10.65), Inches(0.46), font_size=9, color=LIGHTGRAY)

# RAG-specific issue
add_rect(s3, Inches(0.25), Inches(4.3), Inches(12.85), Inches(2.95),
         fill_color=DARKGRAY, line_color=ACCENT, line_width=Pt(1.5))
add_rect(s3, Inches(0.25), Inches(4.3), Inches(12.85), Inches(0.07), fill_color=ACCENT)
add_text_box(s3, "Issue 3: Retrieval Confidence — Averaging Reranker Scores Gave False Negatives (RAG Pipeline)",
             Inches(0.37), Inches(4.39), Inches(12.5), Inches(0.3), font_size=12, bold=True, color=ACCENT2)

rag_cols = [
    ("SYMPTOM", "Retrieving a highly relevant document alongside many irrelevant ones produced strongly negative average reranker score, triggering unnecessary corrective retrieval even when the top result was correct.", RED_SOFT),
    ("ROOT CAUSE", "Averaging all candidate scores (cross-encoder/ms-marco-MiniLM) is misleading — scores range from +5 (relevant) to -11 (irrelevant). Averaging dilutes the signal of one strong positive result.", AMBER),
    ("RESOLUTION", "Switched confidence metric to TOP SCORE + SCORE MARGIN. High top score (>2.0) + large margin (>5.0) = HIGH confidence. Added retrieval agreement signal: if both BM25 and vector independently retrieved the document at rank ≤10, MEDIUM is upgraded to HIGH.", GREEN),
    ("LEARNING", "In cross-encoder reranking, top score and inter-result margin are the reliable signals — not aggregate statistics. Retrieval agreement between BM25 and vector search is a strong signal of document relevance.", ACCENT),
]

rcol_w = Inches(3.0)
rcol_x = Inches(0.35)
for j, (label, text, lc) in enumerate(rag_cols):
    cx = rcol_x + j * (rcol_w + Inches(0.12))
    add_text_box(s3, label, cx, Inches(4.73), rcol_w, Inches(0.22), font_size=7.5, bold=True, color=lc)
    add_text_box(s3, text, cx, Inches(4.95), rcol_w, Inches(1.2), font_size=8.5, color=LIGHTGRAY)

notes_3 = s3.notes_slide.notes_text_frame
notes_3.text = (
    "SLIDE 3 — Production Issues & Resolutions\n\n"
    "Issue 1 (CORS): The root cause was a three-layer problem — "
    "BaseHTTPMiddleware silently consumed CORS headers, and a ghost uvicorn process was still bound to port 8000. "
    "Resolution: Pure-ASGI ExtensionCORSMiddleware shim + kill ghost processes via WMI.\n\n"
    "Issue 2 (Context Extraction): chrome.tabs.sendMessage requires content scripts to be pre-registered. "
    "Replaced with chrome.scripting.executeScript which injects extraction logic on demand, "
    "eliminating content script lifecycle dependency.\n\n"
    "Issue 3 (RAG Confidence): Average reranker score over all candidates was misleading. "
    "Switched to top-score + score-margin + retrieval agreement composite signal, "
    "with upgrade/downgrade rules based on BM25 and vector search agreement."
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — AI Agent Development
# ══════════════════════════════════════════════════════════════════════════════

s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, BG)
add_rect(s4, 0, 0, Inches(0.06), SLIDE_H, fill_color=AMBER)

add_rect(s4, Inches(12.6), Inches(0.12), Inches(0.55), Inches(0.3), fill_color=AMBER)
add_text_box(s4, "04 / 05", Inches(12.58), Inches(0.1), Inches(0.6), Inches(0.35),
             font_size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)

add_text_box(s4, "AI Agent Development", Inches(0.5), Inches(0.12), Inches(11), Inches(0.55),
             font_size=32, bold=True, color=WHITE)
add_text_box(s4, "Agentic retrieval design, orchestration, and autonomous workflow",
             Inches(0.5), Inches(0.62), Inches(11), Inches(0.28), font_size=12, color=GRAY)
add_divider(s4, Inches(0.95), color=AMBER)

# ── Agentic RAG loop diagram ─────────────────────────────────────────────────
section_header(s4, "Agentic / Corrective RAG Loop", Inches(1.05))

loop_steps = [
    ("User Query",         WHITE,   DARKGRAY),
    ("Query Transformer",  ACCENT2, DARKGRAY),
    ("Hybrid Retrieval\nBM25 + Vector",  ACCENT, DARKGRAY),
    ("RRF Fusion",         ACCENT,  DARKGRAY),
    ("Cross-Encoder\nReranker",  ACCENT, DARKGRAY),
    ("Confidence\nGrader",       AMBER,  DARKGRAY),
    ("Answer\nExtractor",        GREEN,  DARKGRAY),
    ("Final Answer\n+ Citation", GREEN,  RGBColor(0x14,0x28,0x1a)),
]

box_w2 = Inches(1.35)
box_h2 = Inches(0.72)
gap2   = Inches(0.1)
sx2    = Inches(0.25)
sy2    = Inches(1.38)
total_w = len(loop_steps) * (box_w2 + gap2) - gap2

for i, (label, tc, fc) in enumerate(loop_steps):
    bx = sx2 + i * (box_w2 + gap2)
    border_c = tc if tc != WHITE else GRAY
    add_rect(s4, bx, sy2, box_w2, box_h2, fill_color=fc, line_color=border_c, line_width=Pt(1.5))
    add_text_box(s4, label, bx + Inches(0.05), sy2 + Inches(0.06),
                 box_w2 - Inches(0.08), box_h2 - Inches(0.06), font_size=9, bold=True,
                 color=tc, align=PP_ALIGN.CENTER)
    if i < len(loop_steps) - 1:
        add_text_box(s4, "→", bx + box_w2 + Inches(0.01), sy2 + Inches(0.22),
                     gap2, Inches(0.3), font_size=12, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)

# Corrective loop back-arrow label
add_rect(s4, Inches(7.5), Inches(2.12), Inches(3.8), Inches(0.38),
         fill_color=RGBColor(0x2a, 0x18, 0x05), line_color=AMBER, line_width=Pt(1))
add_text_box(s4, "↩  Corrective Retrieval: LOW/MEDIUM confidence → synonym expansion or query simplification → re-retrieve",
             Inches(7.55), Inches(2.14), Inches(3.7), Inches(0.36), font_size=8, color=AMBER)

# ── Agentic capabilities table ───────────────────────────────────────────────
section_header(s4, "Implemented Agentic Capabilities", Inches(2.65))

cap_rows = [
    ("Capability", "Implementation", "Code Location", "Status"),
    ("Query Decomposition", "Multi-subquery splitting on 'and/or/;' — each subquery runs independently", "retrieval/query_transform.py · transform_query_full()", "✅ Done"),
    ("Typo & Acronym Correction", "Rule-based typo map + acronym expansion (2FA→two-factor auth, WFH→work from home)", "retrieval/query_transform.py · _correct_typos(), _expand_acronyms()", "✅ Done"),
    ("Confidence Grading", "Composite HIGH/MEDIUM/LOW — top score + margin + BM25↔vector agreement", "retrieval/retrieval_grader.py · grade_retrieval()", "✅ Done"),
    ("Corrective Retrieval", "On LOW/MEDIUM: tries synonym expansion then query simplification; adopts first that improves score", "ui/streamlit_app.py · _CORRECTIVE_STRATEGIES loop", "✅ Done"),
    ("Hybrid Retrieval", "BM25 Okapi + ChromaDB vector search in parallel → RRF fusion (k=60)", "retrieval/bm25_retriever.py + vector_retriever.py + rrf.py", "✅ Done"),
    ("Fallback Strategy", "If no confident span: return full top chunk rather than silence (HIGH/MEDIUM only)", "generation/extractive_qa.py · fallback_to_context=True", "✅ Done"),
    ("Context Capture Agent", "Browser extension as capture agent: auto-scroll → extract → idempotency → queue → backend", "src/background.ts · handleCaptureRequest()", "✅ Done"),
    ("Durable Task Queue", "CaptureQueue with exponential backoff, max 5 attempts, TTL cleanup, chrome.storage.local", "src/core/CaptureQueue.ts", "✅ Done"),
]

col_widths = [Inches(2.1), Inches(3.8), Inches(3.8), Inches(1.8)]
header_colors = [DARKGRAY] * 4
row_y = Inches(2.95)
col_xs = [Inches(0.25), Inches(2.35), Inches(6.15), Inches(9.95)]

for ri, row in enumerate(cap_rows):
    for ci, (text, cx, cw) in enumerate(zip(row, col_xs, col_widths)):
        is_header = ri == 0
        fc = ACCENT if is_header else DARKGRAY
        bc = ACCENT if is_header else RGBColor(0x2a, 0x30, 0x4e)
        add_rect(s4, cx, row_y, cw, Inches(0.36),
                 fill_color=fc, line_color=bc, line_width=Pt(0.5))
        tc2 = WHITE if is_header else (GREEN if "✅" in text else LIGHTGRAY)
        add_text_box(s4, text, cx + Inches(0.06), row_y + Inches(0.04),
                     cw - Inches(0.1), Inches(0.3), font_size=8.5 if not is_header else 9,
                     color=tc2, bold=is_header)
    row_y += Inches(0.36)

notes_4 = s4.notes_slide.notes_text_frame
notes_4.text = (
    "SLIDE 4 — AI Agent Development\n\n"
    "While this project does not use LangChain, LangGraph, CrewAI or MCP, "
    "it implements several agentic design patterns from first principles:\n\n"
    "1. Query Decomposition: Multi-part queries are split into independent subqueries, each run through the full pipeline\n"
    "2. Corrective Retrieval: A feedback loop where the system grades its own retrieval confidence and retries with a modified query strategy\n"
    "3. Confidence Grading: The system makes an autonomous decision about whether to return an answer, trigger corrective retrieval, or return 'Insufficient evidence'\n"
    "4. Chrome Extension as Capture Agent: Autonomously detects AI platform, auto-scrolls for lazy-loaded content, extracts structured conversation, "
    "generates idempotency key, queues with retry — all without user intervention\n"
    "5. Durable Task Queue: CaptureQueue persists to chrome.storage.local and survives service worker restarts\n\n"
    "Next step: integrate LangGraph to formalize the corrective retrieval loop as a graph with conditional edges, "
    "and wire the Advanced RAG pipeline into the FastAPI backend as a /query endpoint."
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Additional Work
# ══════════════════════════════════════════════════════════════════════════════

s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, BG)
add_rect(s5, 0, 0, Inches(0.06), SLIDE_H, fill_color=GREEN)

add_rect(s5, Inches(12.6), Inches(0.12), Inches(0.55), Inches(0.3), fill_color=GREEN)
add_text_box(s5, "05 / 05", Inches(12.58), Inches(0.1), Inches(0.6), Inches(0.35),
             font_size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)

add_text_box(s5, "Additional Work", Inches(0.5), Inches(0.12), Inches(11), Inches(0.55),
             font_size=32, bold=True, color=WHITE)
add_text_box(s5, "Engineering decisions, production features & roadmap",
             Inches(0.5), Inches(0.62), Inches(11), Inches(0.28), font_size=12, color=GRAY)
add_divider(s5, Inches(0.95), color=GREEN)

# ── Left column: Production features ─────────────────────────────────────────
section_header(s5, "Production-Grade Engineering Features", Inches(1.05))

prod_features = [
    ("Multi-Stage Docker Build",
     "python:3.12-slim builder → runtime. Non-root appuser (UID 1001). curl health-check every 30s. uvloop for async performance.",
     GREEN),
    ("Alembic Migration Pipeline",
     "5 versioned migrations (0001→0005). Adds idempotency_key UNIQUE partial index, title, platform, chat_url, messages_count columns to contexts.",
     ACCENT2),
    ("Repository + Service Pattern",
     "ContextRepository, SessionRepository, ProjectRepository — clean separation of SQL from business logic. get_counts() via correlated subqueries.",
     ACCENT),
    ("Structured Logging",
     "structlog with JSON output. Every request logs method, path, status_code, duration_ms, request_id. Context operations log platform, messages count.",
     AMBER),
    ("Test Suite",
     "pytest-asyncio: test_health, test_projects, test_sessions, test_contexts, test_base_repository, test_project_service. httpx TestClient.",
     RED_SOFT),
    ("Idempotent Capture",
     "SHA-256 of (projectId + URL + minute-bucket). Backend checks idempotency_key before insert — duplicate captures return HTTP 200 with original.",
     GREEN),
]

feat_y = Inches(1.35)
for title, body, color in prod_features:
    add_rect(s5, Inches(0.25), feat_y, Inches(6.0), Inches(0.72),
             fill_color=DARKGRAY, line_color=color, line_width=Pt(1.2))
    add_rect(s5, Inches(0.25), feat_y, Inches(0.06), Inches(0.72), fill_color=color)
    add_text_box(s5, title, Inches(0.4), feat_y + Inches(0.05),
                 Inches(5.7), Inches(0.25), font_size=10, bold=True, color=color)
    add_text_box(s5, body, Inches(0.4), feat_y + Inches(0.27),
                 Inches(5.7), Inches(0.44), font_size=8.5, color=LIGHTGRAY)
    feat_y += Inches(0.78)

# ── Right column: Extension architecture highlights + roadmap ─────────────────
section_header(s5, "Extension Architecture Highlights", Inches(1.05))

ext_feats = [
    ("Vite IIFE Build (not ESM)",
     "Chrome extension resources served as application/octet-stream — ESM type=module causes MIME error. Switched to two IIFE Vite configs: vite.popup.config.ts + vite.sidepanel.config.ts.",
     ACCENT),
    ("Zustand + chrome.storage Adapter",
     "Custom chromeStorage adapter persists Zustand store subset (activeProjectId, projects) across popup opens. No localStorage (unavailable in SW).",
     ACCENT2),
    ("CaptureQueue — Durable Retry",
     "Entries: id, projectId, payload, status, attempts, lastAttempt. Exponential backoff, max 5 attempts. TTL: done=1 day, failed=7 days. Survives SW restarts.",
     AMBER),
    ("Platform Auto-Detection",
     "Injected script reads document.location.hostname to detect chatgpt.com, claude.ai, gemini.google.com, perplexity.ai. Falls back to document.body.innerText (50k chars).",
     GREEN),
    ("DiagnosticsBar",
     "Live backend health dot, project count badge, last sync timestamp, retry button on error — always visible in popup header.",
     GRAY),
]

rfeat_y = Inches(1.35)
for title, body, color in ext_feats:
    add_rect(s5, Inches(6.6), rfeat_y, Inches(6.5), Inches(0.85),
             fill_color=DARKGRAY, line_color=color, line_width=Pt(1.2))
    add_rect(s5, Inches(6.6), rfeat_y, Inches(0.06), Inches(0.85), fill_color=color)
    add_text_box(s5, title, Inches(6.75), rfeat_y + Inches(0.05),
                 Inches(6.2), Inches(0.25), font_size=10, bold=True, color=color)
    add_text_box(s5, body, Inches(6.75), rfeat_y + Inches(0.27),
                 Inches(6.2), Inches(0.56), font_size=8.5, color=LIGHTGRAY)
    rfeat_y += Inches(0.91)

# ── Roadmap bar ───────────────────────────────────────────────────────────────
add_rect(s5, Inches(0.25), Inches(6.12), Inches(12.85), Inches(1.2),
         fill_color=RGBColor(0x0a, 0x14, 0x28), line_color=ACCENT, line_width=Pt(1))
add_text_box(s5, "ROADMAP", Inches(0.35), Inches(6.14), Inches(1.2), Inches(0.3),
             font_size=9, bold=True, color=ACCENT)

roadmap = [
    ("Connect RAG → Backend", "POST /api/v1/projects/{id}/query endpoint; route captured contexts through RAG pipeline"),
    ("LangGraph Corrective Loop", "Formalize corrective retrieval as a graph with conditional edges (grade → branch → retry)"),
    ("Semantic Search", "Embed stored contexts with all-MiniLM-L6-v2; expose similarity search across projects"),
    ("LLM Answer Generation", "Add GPT-4o / Claude 3.5 generative answer layer on top of extractive QA"),
    ("Dashboard Analytics", "Token usage, capture frequency, top topics per project via JSONB aggregation"),
]

rmap_x = Inches(1.6)
for title, body in roadmap:
    rw = Inches(2.2)
    add_rect(s5, rmap_x, Inches(6.17), rw, Inches(1.05),
             fill_color=DARKGRAY, line_color=ACCENT2, line_width=Pt(0.8))
    add_text_box(s5, title, rmap_x + Inches(0.07), Inches(6.18),
                 rw - Inches(0.1), Inches(0.28), font_size=8.5, bold=True, color=ACCENT2)
    add_text_box(s5, body, rmap_x + Inches(0.07), Inches(6.45),
                 rw - Inches(0.1), Inches(0.7), font_size=7.5, color=LIGHTGRAY)
    rmap_x += rw + Inches(0.08)

notes_5 = s5.notes_slide.notes_text_frame
notes_5.text = (
    "SLIDE 5 — Additional Work\n\n"
    "Production features actually implemented:\n"
    "- Multi-stage Docker build with non-root user and health check\n"
    "- Alembic migration pipeline with 5 ordered versions\n"
    "- Repository + Service pattern for clean separation of SQL and business logic\n"
    "- structlog structured logging (JSON) on every request\n"
    "- pytest-asyncio test suite covering all API routes and core services\n"
    "- SHA-256 idempotent capture preventing duplicate context storage\n\n"
    "Chrome Extension engineering decisions:\n"
    "- Switched from Vite ESM to IIFE format to resolve MIME type errors in Chrome extension loading\n"
    "- Custom Zustand adapter for chrome.storage.local persistence across popup opens\n"
    "- Durable CaptureQueue survives Service Worker restarts with exponential backoff\n\n"
    "Roadmap: connect RAG engine to backend as /query endpoint, "
    "formalize corrective retrieval with LangGraph, "
    "add generative answer layer (GPT-4o or Claude)."
)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = r"c:\Users\kale_p\Music\Connecting dots\GenAI_Project_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
